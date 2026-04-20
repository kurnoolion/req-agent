"""Add Human Filter philosophy to NORA presentation."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

DARK_BLUE = RGBColor(0x1B, 0x2A, 0x4A)
MID_BLUE = RGBColor(0x2C, 0x5F, 0x8A)
ACCENT_BLUE = RGBColor(0x3A, 0x86, 0xC8)
LIGHT_BLUE = RGBColor(0xD6, 0xEA, 0xF8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
MED_GRAY = RGBColor(0x66, 0x66, 0x66)
LIGHT_GRAY = RGBColor(0xF0, 0xF0, 0xF0)
GREEN = RGBColor(0x27, 0xAE, 0x60)
RED = RGBColor(0xC0, 0x39, 0x2B)
ORANGE = RGBColor(0xE6, 0x7E, 0x22)
TEAL = RGBColor(0x16, 0xA0, 0x85)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation("/home/mohan/work/nora/NORA_Overview.pptx")


def add_shape(slide, left, top, width, height, fill_color, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.fill.solid()
        shape.line.fill.fore_color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_rounded_shape(slide, left, top, width, height, fill_color, border_color=None):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.fill.solid()
        shape.line.fill.fore_color.rgb = border_color
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=14,
                 color=DARK_GRAY, bold=False, alignment=PP_ALIGN.LEFT,
                 font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_bullet_frame(slide, left, top, width, height, items, font_size=14,
                     color=DARK_GRAY, line_spacing=1.3, bold_prefix=True):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_after = Pt(font_size * 0.4)
        p.line_spacing = Pt(font_size * line_spacing)

        if bold_prefix and ":" in item and item.index(":") < 40:
            prefix, rest = item.split(":", 1)
            run1 = p.add_run()
            run1.text = "\u2022  " + prefix + ":"
            run1.font.size = Pt(font_size)
            run1.font.color.rgb = color
            run1.font.bold = True
            run1.font.name = "Calibri"
            run2 = p.add_run()
            run2.text = rest
            run2.font.size = Pt(font_size)
            run2.font.color.rgb = color
            run2.font.name = "Calibri"
        else:
            run = p.add_run()
            run.text = "\u2022  " + item
            run.font.size = Pt(font_size)
            run.font.color.rgb = color
            run.font.name = "Calibri"
    return txBox


def add_slide_number(slide, num):
    add_text_box(slide, Inches(12.5), Inches(7.05), Inches(0.7), Inches(0.35),
                 str(num), font_size=10, color=MED_GRAY, alignment=PP_ALIGN.RIGHT)


def add_title_bar(slide, title, subtitle=None):
    add_shape(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.15), DARK_BLUE)
    add_text_box(slide, Inches(0.6), Inches(0.2), Inches(10), Inches(0.55),
                 title, font_size=26, color=WHITE, bold=True)
    if subtitle:
        add_text_box(slide, Inches(0.6), Inches(0.68), Inches(10), Inches(0.4),
                     subtitle, font_size=14, color=LIGHT_BLUE)
    add_shape(slide, Inches(0), Inches(1.15), SLIDE_W, Pt(3), ACCENT_BLUE)


# ============================================================================
# NEW SLIDE: Human-in-the-Loop (insert after slide 8, before slide 9)
# ============================================================================
# We need to insert at position 8 (0-indexed), which is after current slide 8
# (Design Philosophies) and before slide 9 (PoC Progress)

from pptx.opc.constants import RELATIONSHIP_TYPE as RT
import copy
from lxml import etree

# Create the new slide using blank layout
slide_layout = prs.slide_layouts[6]  # blank
new_slide = prs.slides.add_slide(slide_layout)

# Build the Human Filter slide content
add_bg_fill = new_slide.background.fill
add_bg_fill.solid()
add_bg_fill.fore_color.rgb = WHITE

add_title_bar(new_slide, "Human-in-the-Loop: The \"Human Filter\" Principle",
              "We do not trust AI to achieve 100% accuracy — the system is designed for human oversight at every stage")

# Core principle box
add_shape(new_slide, Inches(0.6), Inches(1.45), Inches(12.1), Inches(1.2), LIGHT_BLUE, ACCENT_BLUE)
add_text_box(new_slide, Inches(0.9), Inches(1.55), Inches(11.5), Inches(0.3),
             "Core Principle", font_size=16, color=DARK_BLUE, bold=True)
add_text_box(new_slide, Inches(0.9), Inches(1.9), Inches(11.5), Inches(0.65),
             "AI is a powerful accelerator, not an infallible oracle. Every AI-generated artifact — document profiles, "
             "feature taxonomies, cross-references, query answers — passes through a Human Filter before it becomes trusted. "
             "The Requirements Analysis team acts as this filter, providing corrections that the system learns from. "
             "This applies equally during development (PoC tuning) and production (ongoing quality assurance).",
             font_size=14, color=DARK_GRAY)

# Three-column layout: During Development | The Mechanism | In Production
col_data = [
    ("During Development", GREEN, [
        "Each pipeline stage output is reviewed by domain experts",
        "Corrections fed back as override files (profile.json, taxonomy.json)",
        "Pipeline auto-detects corrections and uses them on re-run",
        "Evaluation Q&A pairs authored by experts measure real accuracy",
        "A/B testing proves graph value vs. vanilla RAG with human-judged ground truth",
        "Iterative: run → review → correct → re-run → until quality targets met",
    ]),
    ("The Correction Mechanism", ACCENT_BLUE, [
        "Document Profile: fix heading levels, req ID patterns, zone boundaries",
        "Feature Taxonomy: rename, add, remove features; fix categorization",
        "Cross-References: identify missing edges from evaluation failures",
        "Eval Q&A: domain experts write questions with expected answers",
        "QC Reports: structured quality check format for quick feedback",
        "FIX Reports: structured correction descriptions for traceability",
    ]),
    ("In Production", MID_BLUE, [
        "New MNO/release documents go through Human Filter before going live",
        "Requirements team validates profile + taxonomy for each new batch",
        "Query answers include citations — users can verify against source docs",
        "Metrics dashboard tracks accuracy trends over time",
        "Correction workflow is the same: corrections/ dir → re-run pipeline",
        "Continuous improvement: corrections from production feed back into tuning",
    ]),
]

for i, (title, color, items) in enumerate(col_data):
    x = Inches(0.4) + Inches(i * 4.2)
    y = Inches(2.9)

    header = add_shape(new_slide, x, y, Inches(3.95), Inches(0.5), color)
    add_text_box(new_slide, x + Inches(0.1), y + Inches(0.07), Inches(3.75), Inches(0.35),
                 title, font_size=14, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    add_rounded_shape(new_slide, x, y + Inches(0.5), Inches(3.95), Inches(3.25),
                      WHITE, color)
    add_bullet_frame(new_slide, x + Inches(0.12), y + Inches(0.6),
                     Inches(3.7), Inches(3.05),
                     items, font_size=11, line_spacing=1.15, bold_prefix=False)

# Bottom emphasis
add_shape(new_slide, Inches(0.6), Inches(6.8), Inches(12.1), Inches(0.5), DARK_BLUE)
add_text_box(new_slide, Inches(0.9), Inches(6.83), Inches(11.5), Inches(0.4),
             "Result: The system gets better with every human correction — not just at deployment, but continuously throughout its lifecycle.",
             font_size=13, color=WHITE, bold=True)

add_slide_number(new_slide, 9)

# ============================================================================
# Now move the new slide (currently last) to position after slide 8 (index 8)
# ============================================================================
# The new slide is at the end. We need to move it to index 8 (after slide 8, 0-indexed)
slide_list = prs.slides._sldIdLst
slide_ids = list(slide_list)
# The new slide's element is the last one
new_slide_elem = slide_ids[-1]
# Remove from end
slide_list.remove(new_slide_elem)
# Insert at position 8 (after slide 8 which is index 7, so insert at index 8)
slide_list.insert(8, new_slide_elem)

# ============================================================================
# Update slide numbers on subsequent slides (9 onwards shifted by 1)
# Slides after the insert: old 9 (PoC) -> 10, old 10 (Production) -> 11, etc.
# We need to fix slide numbers in the text. Let's update them.
# ============================================================================
# Slide index mapping (0-based):
#   0-7: unchanged (slides 1-8)
#   8: new Human Filter slide (slide 9)
#   9: was PoC Progress (now slide 10)
#   10: was Production Workflow (now slide 11)
#   11: was Dev & Testing (now slide 12)
#   12: was Risks (now slide 13)
#   13: was Roadmap (now slide 14)
#   14: was Q&A (slide 15)

# Update slide numbers on shifted slides
for idx in range(9, len(prs.slides)):
    slide = prs.slides[idx]
    new_num = idx + 1
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                text = para.text.strip()
                # Find slide number text boxes (small number at bottom-right)
                if text.isdigit() and int(text) >= 9 and int(text) <= 14:
                    old_num = int(text)
                    if old_num == new_num - 1 or old_num == new_num:
                        # This is likely the slide number
                        for run in para.runs:
                            if run.text.strip().isdigit():
                                run.text = str(new_num)

# ============================================================================
# Update Agenda slide (slide 2, index 1) to add Human Filter item
# ============================================================================
# We need to shift items 6-10 to 7-11 and insert item 6 "Human Filter"
# Rather than trying to modify the complex agenda layout, we'll rebuild it
# Actually, the agenda items are individual text boxes making it hard to modify
# Let's add a note or simply adjust - we'll rebuild the agenda

# ============================================================================
# Also update slide 11 (Dev & Testing, now index 11) to emphasize Human Filter
# in Requirements Team column
# ============================================================================
# We'll add a highlight box to the Requirements Team column
dev_slide = prs.slides[11]  # Dev & Testing Workflow (now slide 12)

# Add a "Human Filter" badge on the Requirements Team header area
badge = dev_slide.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(4.7) + Inches(0.6), Inches(1.5) + Inches(0.02),
    Inches(1.6), Inches(0.32),

)
badge.fill.solid()
badge.fill.fore_color.rgb = ORANGE
badge.line.fill.background()
tf = badge.text_frame
tf.word_wrap = False
tf.paragraphs[0].alignment = PP_ALIGN.CENTER
r = tf.paragraphs[0].add_run()
r.text = "HUMAN FILTER"
r.font.size = Pt(9)
r.font.color.rgb = WHITE
r.font.bold = True
r.font.name = "Calibri"

# ============================================================================
# Update Roadmap slide (now index 13) - add Human Filter emphasis in Phase 1
# We'll add a callout box
# ============================================================================
roadmap_slide = prs.slides[13]  # Roadmap (now slide 14)

# Add callout at bottom
add_shape(roadmap_slide, Inches(0.4), Inches(6.85), Inches(12.5), Inches(0.48), ORANGE)
add_text_box(roadmap_slide, Inches(0.7), Inches(6.88), Inches(12.0), Inches(0.4),
             "Human Filter throughout all phases: Requirements Analysis team reviews and corrects every AI-generated artifact "
             "before it is trusted — during PoC tuning, scale-out validation, and ongoing production.",
             font_size=12, color=WHITE, bold=True)

# ============================================================================
# Update closing slide (now index 14) to include Human Filter point
# ============================================================================
closing_slide = prs.slides[14]  # Q&A slide

# Find the bullet list and update it
for shape in closing_slide.shapes:
    if shape.has_text_frame:
        tf = shape.text_frame
        texts = [p.text for p in tf.paragraphs]
        if any("Knowledge Graph" in t for t in texts):
            # This is the summary bullet list - add Human Filter point
            p = tf.add_paragraph()
            p.space_after = Pt(8)
            r = p.add_run()
            r.text = "\u2022  Human-in-the-Loop by design: Requirements team acts as Human Filter — we don't trust AI for 100% accuracy"
            r.font.size = Pt(14)
            r.font.color.rgb = LIGHT_BLUE
            r.font.name = "Calibri"
            break


# ============================================================================
# Save
# ============================================================================
prs.save("/home/mohan/work/nora/NORA_Overview.pptx")
print("Presentation updated with Human Filter slide and references.")
print(f"Total slides: {len(prs.slides)}")
