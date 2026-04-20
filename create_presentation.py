"""Generate NORA leadership presentation."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# -- Brand colors --
DARK_BLUE = RGBColor(0x1B, 0x2A, 0x4A)
MID_BLUE = RGBColor(0x2C, 0x5F, 0x8A)
ACCENT_BLUE = RGBColor(0x3A, 0x86, 0xC8)
LIGHT_BLUE = RGBColor(0xD6, 0xEA, 0xF8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
MED_GRAY = RGBColor(0x66, 0x66, 0x66)
LIGHT_GRAY = RGBColor(0xF0, 0xF0, 0xF0)
GREEN = RGBColor(0x27, 0xAE, 0x60)
RED = RGBColor(0xC0, 0x39, 0x2B)
ORANGE = RGBColor(0xE6, 0x7E, 0x22)
TEAL = RGBColor(0x16, 0xA0, 0x85)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def add_bg(slide, color=WHITE):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, fill_color, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.fill.solid()
        shape.line.fill.fore_color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_rounded_shape(slide, left, top, width, height, fill_color, border_color=None):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.fill.solid()
        shape.line.fill.fore_color.rgb = border_color
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=14,
                 color=DARK_GRAY, bold=False, alignment=PP_ALIGN.LEFT,
                 font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_bullet_frame(slide, left, top, width, height, items, font_size=14,
                     color=DARK_GRAY, line_spacing=1.3, bold_prefix=True):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_after = Pt(font_size * 0.4)
        p.line_spacing = Pt(font_size * line_spacing)

        if bold_prefix and ":" in item and item.index(":") < 40:
            prefix, rest = item.split(":", 1)
            run1 = p.add_run()
            run1.text = "\u2022  " + prefix + ":"
            run1.font.size = Pt(font_size)
            run1.font.color.rgb = color
            run1.font.bold = True
            run1.font.name = "Calibri"
            run2 = p.add_run()
            run2.text = rest
            run2.font.size = Pt(font_size)
            run2.font.color.rgb = color
            run2.font.name = "Calibri"
        else:
            run = p.add_run()
            run.text = "\u2022  " + item
            run.font.size = Pt(font_size)
            run.font.color.rgb = color
            run.font.name = "Calibri"
    return txBox


def add_slide_number(slide, num):
    add_text_box(slide, Inches(12.5), Inches(7.05), Inches(0.7), Inches(0.35),
                 str(num), font_size=10, color=MED_GRAY, alignment=PP_ALIGN.RIGHT)


def add_title_bar(slide, title, subtitle=None):
    add_shape(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.15), DARK_BLUE)
    add_text_box(slide, Inches(0.6), Inches(0.2), Inches(10), Inches(0.55),
                 title, font_size=26, color=WHITE, bold=True)
    if subtitle:
        add_text_box(slide, Inches(0.6), Inches(0.68), Inches(10), Inches(0.4),
                     subtitle, font_size=14, color=LIGHT_BLUE)
    add_shape(slide, Inches(0), Inches(1.15), SLIDE_W, Pt(3), ACCENT_BLUE)


def add_card(slide, left, top, width, height, title, body_items,
             title_color=DARK_BLUE, fill=WHITE, border=ACCENT_BLUE,
             font_size=12, title_size=14):
    shape = add_rounded_shape(slide, left, top, width, height, fill, border)
    add_text_box(slide, left + Inches(0.15), top + Inches(0.08),
                 width - Inches(0.3), Inches(0.35),
                 title, font_size=title_size, color=title_color, bold=True)
    add_bullet_frame(slide, left + Inches(0.15), top + Inches(0.42),
                     width - Inches(0.3), height - Inches(0.5),
                     body_items, font_size=font_size, line_spacing=1.2)
    return shape


# ============================================================================
# SLIDE 1: Title
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide, DARK_BLUE)

add_shape(slide, Inches(0), Inches(2.5), SLIDE_W, Inches(3.2), MID_BLUE)

add_text_box(slide, Inches(1), Inches(0.8), Inches(11), Inches(0.6),
             "NORA", font_size=52, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1), Inches(1.6), Inches(11), Inches(0.5),
             "Network Operator Requirements Analyzer", font_size=28,
             color=LIGHT_BLUE, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1.5), Inches(2.8), Inches(10), Inches(0.5),
             "AI-Powered Intelligent Querying, Cross-Referencing, and Compliance Analysis",
             font_size=20, color=WHITE, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1.5), Inches(3.4), Inches(10), Inches(0.5),
             "for US MNO Device Requirement Specifications",
             font_size=20, color=WHITE, alignment=PP_ALIGN.CENTER)

divider = add_shape(slide, Inches(5.5), Inches(4.2), Inches(2.3), Pt(2), ACCENT_BLUE)

add_text_box(slide, Inches(1), Inches(4.6), Inches(11), Inches(0.4),
             "Knowledge Graph + RAG Hybrid Architecture",
             font_size=18, color=LIGHT_BLUE, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1), Inches(5.8), Inches(11), Inches(0.4),
             "Leadership & Cross-Team Review", font_size=16,
             color=RGBColor(0x88, 0xAA, 0xCC), alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1), Inches(6.2), Inches(11), Inches(0.4),
             "April 2026", font_size=14,
             color=RGBColor(0x88, 0xAA, 0xCC), alignment=PP_ALIGN.CENTER)


# ============================================================================
# SLIDE 2: Agenda
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_title_bar(slide, "Agenda")

agenda_items = [
    ("1", "The Problem", "MNO requirement complexity and current pain points"),
    ("2", "Why Vanilla RAG Falls Short", "Concrete example of retrieval failures"),
    ("3", "Our Approach: KG + RAG Hybrid", "New techniques and why they work"),
    ("4", "Pipeline Architecture", "9-stage ingestion + 6-stage query pipeline"),
    ("5", "Key Design Philosophies", "Modularity, observability, team accessibility"),
    ("6", "PoC Progress", "What we've built and validated so far"),
    ("7", "Production Workflow", "How the system serves end users"),
    ("8", "Development & Testing Workflow", "Team roles and contribution model"),
    ("9", "Risks & Mitigations", "Top challenges and our approach"),
    ("10", "Roadmap", "Phased path from PoC to production"),
]

for i, (num, title, desc) in enumerate(agenda_items):
    y = Inches(1.55) + Inches(i * 0.55)
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.0), y, Inches(0.4), Inches(0.4))
    circle.fill.solid()
    circle.fill.fore_color.rgb = ACCENT_BLUE
    circle.line.fill.background()
    tf = circle.text_frame
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    run = tf.paragraphs[0].add_run()
    run.text = num
    run.font.size = Pt(13)
    run.font.color.rgb = WHITE
    run.font.bold = True
    tf.paragraphs[0].space_before = Pt(0)
    tf.word_wrap = False

    add_text_box(slide, Inches(1.6), y - Pt(1), Inches(3.5), Inches(0.4),
                 title, font_size=15, color=DARK_BLUE, bold=True)
    add_text_box(slide, Inches(5.2), y + Pt(2), Inches(6.5), Inches(0.35),
                 desc, font_size=13, color=MED_GRAY)

add_slide_number(slide, 2)


# ============================================================================
# SLIDE 3: The Problem
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_title_bar(slide, "The Problem", "MNO Device Requirement Specifications are Complex and Interconnected")

# Left side: problem description
add_text_box(slide, Inches(0.6), Inches(1.45), Inches(6), Inches(0.35),
             "What We're Dealing With", font_size=18, color=DARK_BLUE, bold=True)

add_bullet_frame(slide, Inches(0.6), Inches(1.9), Inches(5.8), Inches(2.8), [
    "Multiple MNOs (VZW, ATT, TMO) each with quarterly releases",
    "Hundreds of requirement + test case documents per release (GBs of PDFs, DOCX, XLS)",
    "Deep cross-document dependencies: one capability spans multiple docs",
    "Requirements reference specific 3GPP/GSMA standard sections and releases",
    "Different MNOs have different requirements for the same feature",
    "Version tracking across quarterly releases is critical",
], font_size=14)

# Right side: use cases
add_text_box(slide, Inches(7.0), Inches(1.45), Inches(5.5), Inches(0.35),
             "What Teams Need to Do Today (Manually)", font_size=18, color=DARK_BLUE, bold=True)

add_bullet_frame(slide, Inches(7.0), Inches(1.9), Inches(5.8), Inches(2.8), [
    "Cross-doc Q&A: \"What is the complete T3402 timer behavior?\" (spans 3+ docs)",
    "Cross-MNO comparison: \"How does VZW vs TMO handle IMS registration?\"",
    "Standards tracing: \"How does VZW differ from 3GPP for attach reject?\"",
    "Release diff: \"What changed in VZW eSIM from Oct 2025 to Feb 2026?\"",
    "Req-to-test traceability: \"Which test cases cover this requirement?\"",
    "Compliance analysis: check device implementation against all MNO reqs",
], font_size=14)

# Bottom highlight
add_shape(slide, Inches(0.6), Inches(5.2), Inches(12.1), Inches(1.6), LIGHT_BLUE, ACCENT_BLUE)
add_text_box(slide, Inches(0.9), Inches(5.3), Inches(11.5), Inches(0.35),
             "The Core Challenge", font_size=16, color=DARK_BLUE, bold=True)
add_text_box(slide, Inches(0.9), Inches(5.7), Inches(11.5), Inches(1.0),
             "Engineers spend hours manually cross-referencing documents, tracing requirements to standards, "
             "and comparing across MNOs and releases. The information exists but is scattered across hundreds "
             "of interconnected documents with no automated way to query, cross-reference, or analyze it.",
             font_size=14, color=DARK_GRAY)

add_slide_number(slide, 3)


# ============================================================================
# SLIDE 4: Why Vanilla RAG Fails
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_title_bar(slide, "Why Vanilla RAG Falls Short",
              "Standard vector-based retrieval fundamentally cannot handle this problem")

# Left: limitations
add_text_box(slide, Inches(0.6), Inches(1.5), Inches(5.8), Inches(0.35),
             "6 Fundamental Limitations", font_size=17, color=DARK_BLUE, bold=True)

limitations = [
    "No relationship awareness: treats each chunk independently, can't follow cross-doc dependencies",
    "Undirected retrieval: semantic search over entire corpus, may miss critical related chunks",
    "Destroyed structure: chunking breaks hierarchical parent-child relationships",
    "Poor telecom terminology: standard embeddings not trained on 3GPP acronyms",
    "Missing standards context: \"follow 3GPP TS 24.301 sec 5.5.1\" is incomplete without the standard",
    "No MNO/release scoping: no mechanism to scope results or support cross-MNO comparison",
]

for i, lim in enumerate(limitations):
    y = Inches(2.0) + Inches(i * 0.7)
    x_mark = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.7), y + Pt(2), Inches(0.28), Inches(0.28))
    x_mark.fill.solid()
    x_mark.fill.fore_color.rgb = RED
    x_mark.line.fill.background()
    tf = x_mark.text_frame
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    r = tf.paragraphs[0].add_run()
    r.text = "\u2717"
    r.font.size = Pt(12)
    r.font.color.rgb = WHITE
    r.font.bold = True

    add_text_box(slide, Inches(1.15), y, Inches(5.3), Inches(0.6),
                 lim, font_size=13, color=DARK_GRAY)

# Right: concrete example
add_rounded_shape(slide, Inches(6.8), Inches(1.5), Inches(5.9), Inches(5.3),
                  RGBColor(0xFD, 0xF2, 0xE9), ORANGE)

add_text_box(slide, Inches(7.1), Inches(1.6), Inches(5.3), Inches(0.35),
             "Concrete Example", font_size=17, color=ORANGE, bold=True)

add_text_box(slide, Inches(7.1), Inches(2.0), Inches(5.3), Inches(0.35),
             "Query: \"What is the complete T3402 timer behavior?\"",
             font_size=13, color=DARK_GRAY, bold=True)

add_text_box(slide, Inches(7.1), Inches(2.5), Inches(5.3), Inches(0.3),
             "What Vanilla RAG returns:", font_size=13, color=RED, bold=True)
add_bullet_frame(slide, Inches(7.1), Inches(2.8), Inches(5.3), Inches(1.3), [
    "Chunks from LTEDATARETRY mentioning T3402",
    "Misses: related timer behavior in LTEB13NAC",
    "Misses: 3GPP TS 24.301 sec 5.5.1.2.5 (the standard it defers to)",
    "Misses: parent requirement context and hierarchy",
], font_size=12, color=DARK_GRAY)

add_text_box(slide, Inches(7.1), Inches(4.25), Inches(5.3), Inches(0.3),
             "What the answer actually requires:", font_size=13, color=GREEN, bold=True)
add_bullet_frame(slide, Inches(7.1), Inches(4.55), Inches(5.3), Inches(1.8), [
    "Requirements from LTEDATARETRY (primary definition)",
    "Cross-references from LTEB13NAC (Band 13 specific behavior)",
    "3GPP TS 24.301 section 5.5.1.2.5, Release 10 (deferred standard)",
    "Hierarchical context: parent sections defining retry scope",
    "Feature-level context: related DATA_RETRY feature requirements",
], font_size=12, color=DARK_GRAY)

add_slide_number(slide, 4)


# ============================================================================
# SLIDE 5: Our Approach
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_title_bar(slide, "Our Approach: Knowledge Graph + RAG Hybrid",
              "Use a Knowledge Graph to scope WHERE to look, RAG to rank WHAT's relevant, "
              "hierarchy for structural CONTEXT")

# Core insight box
add_shape(slide, Inches(0.6), Inches(1.45), Inches(12.1), Inches(1.1), LIGHT_BLUE, ACCENT_BLUE)
add_text_box(slide, Inches(0.9), Inches(1.55), Inches(11.5), Inches(0.3),
             "Core Architectural Insight", font_size=15, color=DARK_BLUE, bold=True)
add_text_box(slide, Inches(0.9), Inches(1.9), Inches(11.5), Inches(0.5),
             "Graph scoping identifies WHERE to look (cross-doc, cross-MNO traversal)  \u2192  "
             "Targeted vector RAG ranks WHAT's most relevant within that scope  \u2192  "
             "Requirement hierarchy provides structural CONTEXT for LLM synthesis",
             font_size=14, color=DARK_GRAY)

# Three technique cards
card_data = [
    ("Knowledge Graph", DARK_BLUE, [
        "Unified graph: all MNOs \u00d7 releases \u00d7 docs",
        "8 node types (MNO, Release, Plan, Requirement, Standard, Feature, ...)",
        "15+ edge types (depends_on, references_standard, maps_to, ...)",
        "Natural cross-MNO comparison via shared Feature & Standards nodes",
        "Cross-release diff via version traversal",
    ]),
    ("Targeted RAG", MID_BLUE, [
        "Vector retrieval scoped by graph candidates (not full corpus)",
        "Metadata-filtered by MNO, release, doc_type",
        "Diversity enforcement: minimum chunks per plan",
        "Contextualized chunks with hierarchy path and MNO/release headers",
        "Configurable: embedding model, distance metric, top-k",
    ]),
    ("LLM Synthesis", TEAL, [
        "Context enriched with graph relationships (hierarchy, standards, cross-refs)",
        "Query-type-specific system prompts (comparison, traceability, etc.)",
        "Few-shot citation examples for accurate sourcing",
        "LLM-agnostic: Protocol-based abstraction (swap providers freely)",
        "Citation extraction and fallback for smaller models",
    ]),
]

for i, (title, color, items) in enumerate(card_data):
    x = Inches(0.6) + Inches(i * 4.1)
    add_card(slide, x, Inches(2.85), Inches(3.85), Inches(3.6),
             title, items, title_color=color, border=color, font_size=13)

# Bottom: how example is solved
add_shape(slide, Inches(0.6), Inches(6.65), Inches(12.1), Inches(0.55), DARK_BLUE)
add_text_box(slide, Inches(0.9), Inches(6.7), Inches(11.5), Inches(0.45),
             "T3402 Example Solved:  Graph finds LTEDATARETRY + LTEB13NAC + TS 24.301 nodes  \u2192  "
             "RAG ranks best chunks within those  \u2192  LLM synthesizes complete answer with citations",
             font_size=13, color=WHITE)

add_slide_number(slide, 5)


# ============================================================================
# SLIDE 6: Pipeline Architecture (Ingestion)
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_title_bar(slide, "Ingestion Pipeline", "9-stage document processing: from raw PDFs to queryable knowledge")

stages = [
    ("1", "Extract", "PDF/DOCX/XLS content\n\u2192 normalized IR", RGBColor(0x2E, 0x86, 0xAB)),
    ("2", "Profile", "Derive document\nstructure (LLM-free)", RGBColor(0x2E, 0x86, 0xAB)),
    ("3", "Parse", "IR \u2192 requirement\ntrees (profile-driven)", RGBColor(0x2E, 0x86, 0xAB)),
    ("4", "Test Cases", "Parse test case\ndocuments", RGBColor(0x99, 0x99, 0x99)),
    ("5", "Cross-Refs", "Resolve internal +\ncross-doc + standards", RGBColor(0x1A, 0x7A, 0x5A)),
    ("6", "Taxonomy", "LLM-derived feature\ncategorization", RGBColor(0x1A, 0x7A, 0x5A)),
    ("7", "Standards", "Download + parse\n3GPP spec sections", RGBColor(0x1A, 0x7A, 0x5A)),
    ("8", "Graph", "Build unified\nknowledge graph", RGBColor(0x8E, 0x44, 0xAD)),
    ("9", "Vectors", "Embed chunks \u2192\nvector store", RGBColor(0x8E, 0x44, 0xAD)),
]

for i, (num, name, desc, color) in enumerate(stages):
    x = Inches(0.35) + Inches(i * 1.42)
    y = Inches(1.55)

    box = add_rounded_shape(slide, x, y, Inches(1.28), Inches(1.7), color)
    tf = box.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    r = tf.paragraphs[0].add_run()
    r.text = num
    r.font.size = Pt(22)
    r.font.color.rgb = WHITE
    r.font.bold = True
    r.font.name = "Calibri"

    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = name
    r2.font.size = Pt(13)
    r2.font.color.rgb = WHITE
    r2.font.bold = True
    r2.font.name = "Calibri"

    add_text_box(slide, x - Inches(0.05), y + Inches(1.8), Inches(1.38), Inches(0.8),
                 desc, font_size=10, color=DARK_GRAY, alignment=PP_ALIGN.CENTER)

    if i < len(stages) - 1:
        arrow = slide.shapes.add_shape(
            MSO_SHAPE.RIGHT_ARROW, x + Inches(1.3), y + Inches(0.65),
            Inches(0.14), Inches(0.35)
        )
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = MED_GRAY
        arrow.line.fill.background()

# Phase labels
add_text_box(slide, Inches(0.35), Inches(3.85), Inches(4.3), Inches(0.3),
             "Document Understanding", font_size=11, color=RGBColor(0x2E, 0x86, 0xAB),
             bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(5.3), Inches(3.85), Inches(4.3), Inches(0.3),
             "Relationship Extraction", font_size=11, color=RGBColor(0x1A, 0x7A, 0x5A),
             bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(10.1), Inches(3.85), Inches(2.9), Inches(0.3),
             "Knowledge Storage", font_size=11, color=RGBColor(0x8E, 0x44, 0xAD),
             bold=True, alignment=PP_ALIGN.CENTER)

# Rationale section
add_text_box(slide, Inches(0.6), Inches(4.3), Inches(12), Inches(0.35),
             "Why Each Stage Exists", font_size=16, color=DARK_BLUE, bold=True)

rationale_left = [
    "Extract: Different doc formats (PDF/DOCX/XLS) need unified representation before any analysis",
    "Profile: MNO doc structures vary; auto-derive patterns (LLM-free) instead of hardcoding per-MNO parsers",
    "Parse: Convert flat text into structured requirement trees with hierarchy, req IDs, metadata",
    "Cross-Refs: Requirements reference each other and standards \u2014 these links ARE the knowledge graph edges",
]
rationale_right = [
    "Taxonomy: LLM derives feature categories (\"Data Retry\", \"IMS Registration\") for feature-level queries",
    "Standards: MNO reqs say \"follow 3GPP TS 24.301 sec 5.5.1\" \u2014 incomplete without the actual standard text",
    "Graph: Unified graph enables cross-doc, cross-MNO, cross-release traversal as natural graph operations",
    "Vectors: Embedding-based retrieval within graph-scoped candidates for ranking relevance",
]

add_bullet_frame(slide, Inches(0.6), Inches(4.7), Inches(6.0), Inches(2.6),
                 rationale_left, font_size=12, line_spacing=1.25)
add_bullet_frame(slide, Inches(6.8), Inches(4.7), Inches(6.0), Inches(2.6),
                 rationale_right, font_size=12, line_spacing=1.25)

add_slide_number(slide, 6)


# ============================================================================
# SLIDE 7: Query Pipeline
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_title_bar(slide, "Query Pipeline", "6-stage query processing: from natural language to cited answers")

query_stages = [
    ("1. Query\nAnalysis", "Extract entities,\nfeatures, MNOs,\nintent from NL", ACCENT_BLUE),
    ("2. MNO/Release\nResolution", "Map to available\ngraph scopes\n(latest, all, specific)", ACCENT_BLUE),
    ("3. Graph\nScoping", "Traverse KG to find\ncandidate requirement\nnodes (BFS, 2 hops)", MID_BLUE),
    ("4. Targeted\nRAG", "Vector similarity\nWITHIN candidate set\n(not full corpus)", MID_BLUE),
    ("5. Context\nAssembly", "Enrich chunks with\nhierarchy, standards,\ncross-refs from graph", DARK_BLUE),
    ("6. LLM\nSynthesis", "Generate answer\nwith citations,\nMNO-aware", DARK_BLUE),
]

for i, (name, desc, color) in enumerate(query_stages):
    x = Inches(0.4) + Inches(i * 2.13)
    y = Inches(1.6)

    box = add_rounded_shape(slide, x, y, Inches(1.9), Inches(1.3), color)
    tf = box.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    r = tf.paragraphs[0].add_run()
    r.text = name
    r.font.size = Pt(12)
    r.font.color.rgb = WHITE
    r.font.bold = True
    r.font.name = "Calibri"

    add_text_box(slide, x, y + Inches(1.4), Inches(1.9), Inches(0.9),
                 desc.replace("\n", " "), font_size=11, color=DARK_GRAY,
                 alignment=PP_ALIGN.CENTER)

    if i < len(query_stages) - 1:
        arrow = slide.shapes.add_shape(
            MSO_SHAPE.RIGHT_ARROW,
            x + Inches(1.92), y + Inches(0.45),
            Inches(0.2), Inches(0.35)
        )
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = MED_GRAY
        arrow.line.fill.background()

# Supported query types
add_text_box(slide, Inches(0.6), Inches(4.0), Inches(12), Inches(0.35),
             "Supported Query Types", font_size=16, color=DARK_BLUE, bold=True)

qtypes = [
    ("Single-Doc Q&A", "\"What are the throttling rules in LTEDATARETRY?\""),
    ("Cross-Doc Q&A", "\"What is the complete T3402 timer behavior?\" (spans multiple docs)"),
    ("Cross-MNO Comparison", "\"Compare VZW vs TMO IMS registration requirements\""),
    ("Standards Comparison", "\"How does VZW differ from 3GPP for T3402?\""),
    ("Release Diff", "\"What changed in VZW eSIM from Oct 2025 to Feb 2026?\""),
    ("Traceability", "\"Which test cases cover VZ_REQ_LTEDATARETRY_7748?\""),
]

for i, (qtype, example) in enumerate(qtypes):
    col = i % 2
    row = i // 2
    x = Inches(0.6) + Inches(col * 6.3)
    y = Inches(4.5) + Inches(row * 0.85)

    add_rounded_shape(slide, x, y, Inches(5.9), Inches(0.72), LIGHT_GRAY, ACCENT_BLUE)
    add_text_box(slide, x + Inches(0.15), y + Inches(0.05), Inches(2.2), Inches(0.3),
                 qtype, font_size=12, color=DARK_BLUE, bold=True)
    add_text_box(slide, x + Inches(0.15), y + Inches(0.35), Inches(5.5), Inches(0.3),
                 example, font_size=11, color=MED_GRAY)

add_slide_number(slide, 7)


# ============================================================================
# SLIDE 8: Key Design Philosophies
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_title_bar(slide, "Key Design Philosophies",
              "Principles that guide every architectural decision")

philosophies = [
    ("Modular, Individually Testable Stages",
     "Each pipeline stage is independent with clear inputs/outputs. "
     "Stages can be run, tested, and tuned in isolation. "
     "426 tests verify each stage independently.",
     "\u2699"),
    ("Profile-Driven, Not Hardcoded",
     "DocumentProfiler auto-derives document structure (LLM-free). "
     "Adding a new MNO = profile representative docs, no code changes. "
     "Profiles are human-reviewable and editable JSON.",
     "\U0001F4CB"),
    ("LLM-Agnostic Protocol Design",
     "LLMProvider Protocol (structural typing) \u2014 any class with a matching complete() method works. "
     "Swap between Mock, Ollama, or proprietary LLM with zero code changes. "
     "No vendor lock-in.",
     "\U0001F504"),
    ("Observable and Debuggable",
     "Compact report format (RPT/HW/MDL lines) for cross-environment debugging. "
     "Persistent metrics DB (request timing, LLM stats, resource usage). "
     "Structured error codes (40 codes) for remote troubleshooting.",
     "\U0001F4CA"),
    ("Unified Graph, Not Partitioned",
     "Single graph spans all MNOs \u00d7 releases. Cross-MNO comparison and cross-release diff "
     "are natural traversals, not result merging. Shared Standards and Feature nodes eliminate duplication.",
     "\U0001F310"),
    ("Team-Accessible, Not Expert-Only",
     "Web UI (browser-based) for team members on Windows PCs. "
     "Pipeline submission via forms, real-time job monitoring, shared folder browsing. "
     "No Linux terminal or CLI experience required.",
     "\U0001F465"),
]

for i, (title, desc, icon) in enumerate(philosophies):
    col = i % 2
    row = i // 2
    x = Inches(0.5) + Inches(col * 6.3)
    y = Inches(1.55) + Inches(row * 1.85)

    add_rounded_shape(slide, x, y, Inches(6.0), Inches(1.7), WHITE, ACCENT_BLUE)
    add_text_box(slide, x + Inches(0.7), y + Inches(0.12), Inches(5.1), Inches(0.35),
                 title, font_size=15, color=DARK_BLUE, bold=True)
    add_text_box(slide, x + Inches(0.7), y + Inches(0.5), Inches(5.1), Inches(1.1),
                 desc, font_size=12, color=DARK_GRAY)

    icon_circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, x + Inches(0.12), y + Inches(0.25), Inches(0.45), Inches(0.45)
    )
    icon_circle.fill.solid()
    icon_circle.fill.fore_color.rgb = LIGHT_BLUE
    icon_circle.line.fill.background()

add_slide_number(slide, 8)


# ============================================================================
# SLIDE 9: PoC Progress
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_title_bar(slide, "PoC Progress", "What we've built and validated")

# Left: checklist
add_text_box(slide, Inches(0.6), Inches(1.5), Inches(5.5), Inches(0.35),
             "Implementation Status", font_size=17, color=DARK_BLUE, bold=True)

steps = [
    ("\u2705", "Document Content Extraction (PDF, DOCX, XLS)"),
    ("\u2705", "DocumentProfiler (LLM-free structure analysis)"),
    ("\u2705", "Generic Structural Parser (profile-driven)"),
    ("\u23F3", "Test Case Parser (deferred to after PoC validation)"),
    ("\u2705", "Cross-Reference Resolver (internal + cross-doc + standards)"),
    ("\u2705", "Feature Taxonomy (LLM-derived categorization)"),
    ("\u2705", "Standards Ingestion (3GPP download + parse + extract)"),
    ("\u2705", "Knowledge Graph Construction (1,078 nodes, 11,732 edges)"),
    ("\u2705", "Vector Store (705 contextualized chunks)"),
    ("\u2705", "Query Pipeline (6-stage, 8 query types)"),
    ("\u2705", "Evaluation Framework (18 test questions, A/B comparison)"),
]

for i, (icon, text) in enumerate(steps):
    y = Inches(1.95) + Inches(i * 0.42)
    add_text_box(slide, Inches(0.7), y, Inches(0.35), Inches(0.35),
                 icon, font_size=14, alignment=PP_ALIGN.CENTER)
    color = MED_GRAY if icon == "\u23F3" else DARK_GRAY
    add_text_box(slide, Inches(1.1), y, Inches(5.0), Inches(0.35),
                 text, font_size=13, color=color)

# Right: key metrics + infrastructure
add_text_box(slide, Inches(7.0), Inches(1.5), Inches(5.5), Inches(0.35),
             "Key Metrics (VZW OA, 5 docs)", font_size=17, color=DARK_BLUE, bold=True)

metrics = [
    ("Requirements parsed", "705"),
    ("Graph nodes", "1,078"),
    ("Graph edges", "11,732"),
    ("Features extracted", "16"),
    ("3GPP standards integrated", "11 spec-release pairs"),
    ("Vector chunks", "705"),
    ("Test cases", "426"),
    ("Connected components", "22 (98.1% in largest)"),
]

for i, (label, value) in enumerate(metrics):
    y = Inches(1.95) + Inches(i * 0.4)
    add_text_box(slide, Inches(7.1), y, Inches(3.5), Inches(0.35),
                 label, font_size=13, color=DARK_GRAY)
    add_text_box(slide, Inches(10.5), y, Inches(2.0), Inches(0.35),
                 value, font_size=13, color=ACCENT_BLUE, bold=True)

# Infrastructure built
add_text_box(slide, Inches(7.0), Inches(5.25), Inches(5.5), Inches(0.35),
             "Infrastructure Built", font_size=15, color=DARK_BLUE, bold=True)

add_bullet_frame(slide, Inches(7.0), Inches(5.6), Inches(5.8), Inches(1.5), [
    "Pipeline Runner: 9-stage orchestration, compact reports",
    "Web UI: browser-based access for all team members",
    "Metrics & Observability: persistent DB, resource sampling",
    "Multi-env system: scoped workspaces per team member",
    "Offline install workflow: proxy-restricted environments",
], font_size=12, line_spacing=1.2)

add_slide_number(slide, 9)


# ============================================================================
# SLIDE 10: Production Workflow
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_title_bar(slide, "Production Workflow",
              "How the system serves end users in production")

# Flow diagram using shapes
flow_steps = [
    ("MNO Release\nDropped", "New quarterly docs\nadded to shared folder", RGBColor(0x2E, 0x86, 0xAB)),
    ("Ingestion\nPipeline", "9 stages: extract\nthrough vectorstore", RGBColor(0x1A, 0x7A, 0x5A)),
    ("Unified\nKG + Vectors", "All MNOs, releases,\nstandards indexed", RGBColor(0x8E, 0x44, 0xAD)),
    ("Query\nPipeline", "6 stages: analyze\nthrough synthesize", DARK_BLUE),
    ("Cited\nAnswer", "Answer with req IDs,\nstandards, provenance", RGBColor(0xC0, 0x39, 0x2B)),
]

for i, (title, desc, color) in enumerate(flow_steps):
    x = Inches(0.4) + Inches(i * 2.6)
    y = Inches(1.6)

    box = add_rounded_shape(slide, x, y, Inches(2.2), Inches(1.4), color)
    tf = box.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    r = tf.paragraphs[0].add_run()
    r.text = title
    r.font.size = Pt(14)
    r.font.color.rgb = WHITE
    r.font.bold = True
    r.font.name = "Calibri"

    add_text_box(slide, x, y + Inches(1.5), Inches(2.2), Inches(0.7),
                 desc.replace("\n", " "), font_size=11, color=DARK_GRAY,
                 alignment=PP_ALIGN.CENTER)

    if i < len(flow_steps) - 1:
        arrow = slide.shapes.add_shape(
            MSO_SHAPE.RIGHT_ARROW,
            x + Inches(2.22), y + Inches(0.45),
            Inches(0.35), Inches(0.4)
        )
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = MED_GRAY
        arrow.line.fill.background()

# Use case flows
add_text_box(slide, Inches(0.6), Inches(3.8), Inches(12), Inches(0.35),
             "End-User Interaction Flows", font_size=16, color=DARK_BLUE, bold=True)

flows = [
    ("Requirements Engineer",
     "Queries requirements across docs and MNOs, traces to standards, identifies cross-doc dependencies. "
     "\"What are all IMS registration requirements across VZW and TMO?\""),
    ("QA Engineer",
     "Traces requirements to test cases, identifies coverage gaps, queries test case content. "
     "\"Which test cases cover VZ_REQ_LTEDATARETRY_7748?\""),
    ("Compliance Analyst",
     "Checks device implementation against MNO requirements, generates delta compliance sheets. "
     "\"What changed in VZW eSIM requirements from Oct 2025 to Feb 2026?\""),
    ("Release Manager",
     "Cross-MNO comparison for release planning, identifies MNO-specific customizations over standards. "
     "\"How does VZW differ from 3GPP for attach reject handling?\""),
]

for i, (role, desc) in enumerate(flows):
    y = Inches(4.25) + Inches(i * 0.75)
    add_rounded_shape(slide, Inches(0.6), y, Inches(12.1), Inches(0.65), LIGHT_GRAY)
    add_text_box(slide, Inches(0.8), y + Inches(0.05), Inches(2.2), Inches(0.3),
                 role, font_size=13, color=ACCENT_BLUE, bold=True)
    add_text_box(slide, Inches(3.1), y + Inches(0.05), Inches(9.3), Inches(0.55),
                 desc, font_size=11, color=DARK_GRAY)

add_slide_number(slide, 10)


# ============================================================================
# SLIDE 11: Development & Testing Workflow
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_title_bar(slide, "Development & Testing Workflow",
              "Who contributes what and how")

# Three team columns
teams = [
    ("Orchestrators", "(Architecture & Pipeline)", DARK_BLUE, [
        "Design and implement pipeline stages",
        "Integrate and configure LLM providers",
        "Build evaluation framework and test questions",
        "Create environments for team members",
        "Analyze compact reports and debug remotely",
        "Tune pipeline parameters based on feedback",
        "Maintain Web UI and infrastructure",
    ]),
    ("Requirements Team", "(Domain Expertise)", TEAL, [
        "Run pipeline stages via Web UI or CLI",
        "Review and correct document profiles",
        "Review and correct feature taxonomy",
        "Supply domain-specific eval Q&A pairs (Excel)",
        "Validate cross-reference completeness",
        "Report quality issues (QC template format)",
        "Submit corrections (FIX template format)",
    ]),
    ("QA Team", "(Validation & Coverage)", ORANGE, [
        "Run pipeline on test case documents",
        "Verify requirement-to-test traceability",
        "Evaluate answer quality and citation accuracy",
        "Run A/B comparisons (graph vs pure RAG)",
        "Report evaluation results (compact format)",
        "Identify missing edges and relationships",
        "Validate standards integration accuracy",
    ]),
]

for i, (name, subtitle, color, items) in enumerate(teams):
    x = Inches(0.5) + Inches(i * 4.2)
    y = Inches(1.5)

    header = add_shape(slide, x, y, Inches(3.9), Inches(0.75), color)
    add_text_box(slide, x + Inches(0.1), y + Inches(0.05), Inches(3.7), Inches(0.35),
                 name, font_size=17, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.1), y + Inches(0.38), Inches(3.7), Inches(0.3),
                 subtitle, font_size=11, color=RGBColor(0xDD, 0xDD, 0xDD),
                 alignment=PP_ALIGN.CENTER)

    add_rounded_shape(slide, x, y + Inches(0.75), Inches(3.9), Inches(3.85),
                      WHITE, color)
    add_bullet_frame(slide, x + Inches(0.15), y + Inches(0.9),
                     Inches(3.6), Inches(3.5),
                     items, font_size=12, line_spacing=1.15, bold_prefix=False)

# Bottom: feedback loop
add_shape(slide, Inches(0.5), Inches(6.3), Inches(12.3), Inches(0.85), LIGHT_BLUE, ACCENT_BLUE)
add_text_box(slide, Inches(0.8), Inches(6.35), Inches(11.7), Inches(0.3),
             "Feedback Loop", font_size=14, color=DARK_BLUE, bold=True)
add_text_box(slide, Inches(0.8), Inches(6.65), Inches(11.7), Inches(0.45),
             "Run pipeline  \u2192  Review artifacts  \u2192  Submit corrections (profile, taxonomy, eval Q&A)  "
             "\u2192  Re-run with corrections auto-detected  \u2192  Report results (compact RPT format)  "
             "\u2192  Orchestrators tune parameters  \u2192  Repeat",
             font_size=13, color=DARK_GRAY)

add_slide_number(slide, 11)


# ============================================================================
# SLIDE 12: Risks & Mitigations
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_title_bar(slide, "Risks & Mitigations", "Top challenges and our approach to each")

risks = [
    ("LLM Portability",
     "Patterns that work with open-source Gemma may not transfer to our proprietary LLM",
     "High", "Medium",
     "LLMProvider Protocol abstraction \u2014 swap providers with zero code changes. "
     "Simple prompt patterns, no vendor-specific features. Test portability early in Phase 2."),
    ("Document Extraction Quality",
     "Tables, diagrams, embedded objects may not extract cleanly across all formats",
     "High", "Medium",
     "Multi-library approach (pymupdf + pdfplumber for PDF, python-docx for DOCX). "
     "Prefer DOCX when available. DOC auto-converted via LibreOffice. Quality validated per format."),
    ("Cross-Document Edge Completeness",
     "Implicit dependencies between documents may not be detected, leading to incomplete answers",
     "High", "Medium",
     "Three-layer extraction: deterministic regex + standards reference resolution + LLM-based. "
     "Measure edge recall during PoC. Add missing edges from evaluation failures and team corrections."),
    ("Embedding Quality for Telecom Terms",
     "Standard embedding models not trained on 3GPP acronyms and MNO-specific terminology",
     "Medium", "Medium",
     "Contextualized chunks (hierarchy path + MNO/release headers prepended). "
     "Configurable embedding model \u2014 can swap to domain-adapted models. Evaluate retrieval quality during PoC."),
    ("Feature Taxonomy Accuracy",
     "LLM-derived taxonomy may incorrectly group or miss features, especially cross-MNO alignment",
     "Medium", "Medium",
     "Human review step is mandatory (Requirements Team). Correction workflow: edit taxonomy.json, "
     "place in corrections/, re-run. Start with fewer, broader features and refine iteratively."),
    ("Scale: Unified Graph Growth",
     "Single graph spanning all MNOs \u00d7 releases could become very large in production",
     "Low (PoC)\nHigh (Prod)", "High",
     "PoC validates architecture with single MNO+release. Production uses Neo4j with indexing on mno+release. "
     "Sharding by MNO is a fallback if needed."),
]

headers = ["Risk", "Impact", "Likelihood", "Mitigation"]
col_widths = [Inches(2.5), Inches(0.8), Inches(0.85), Inches(8.0)]
col_starts = [Inches(0.5)]
for w in col_widths[:-1]:
    col_starts.append(col_starts[-1] + w + Inches(0.05))

# Header row
y = Inches(1.5)
for j, (header, cw, cx) in enumerate(zip(headers, col_widths, col_starts)):
    add_shape(slide, cx, y, cw, Inches(0.38), DARK_BLUE)
    add_text_box(slide, cx + Inches(0.08), y + Inches(0.03), cw - Inches(0.15), Inches(0.3),
                 header, font_size=12, color=WHITE, bold=True,
                 alignment=PP_ALIGN.CENTER if j in (1, 2) else PP_ALIGN.LEFT)

# Data rows
for i, (name, _desc, impact, likelihood, mitigation) in enumerate(risks):
    y = Inches(1.95) + Inches(i * 0.85)
    bg_color = LIGHT_GRAY if i % 2 == 0 else WHITE

    for cw, cx in zip(col_widths, col_starts):
        add_shape(slide, cx, y, cw, Inches(0.78), bg_color)

    add_text_box(slide, col_starts[0] + Inches(0.08), y + Inches(0.05),
                 col_widths[0] - Inches(0.1), Inches(0.7),
                 name, font_size=11, color=DARK_BLUE, bold=True)

    severity_color = RED if impact.startswith("High") else ORANGE
    add_text_box(slide, col_starts[1] + Inches(0.08), y + Inches(0.15),
                 col_widths[1] - Inches(0.1), Inches(0.35),
                 impact, font_size=10, color=severity_color, bold=True,
                 alignment=PP_ALIGN.CENTER)

    add_text_box(slide, col_starts[2] + Inches(0.08), y + Inches(0.15),
                 col_widths[2] - Inches(0.1), Inches(0.35),
                 likelihood, font_size=10, color=MED_GRAY,
                 alignment=PP_ALIGN.CENTER)

    add_text_box(slide, col_starts[3] + Inches(0.08), y + Inches(0.05),
                 col_widths[3] - Inches(0.15), Inches(0.7),
                 mitigation, font_size=10, color=DARK_GRAY)

add_slide_number(slide, 12)


# ============================================================================
# SLIDE 13: Roadmap
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_title_bar(slide, "Roadmap", "Phased path from PoC to production")

phases = [
    ("Phase 1: PoC Validation", "Current \u2192 Near-term", GREEN, [
        ("Scope", "1 MNO (VZW), 1 Release (Feb 2026), 5 requirement documents"),
        ("LLM", "Hook up proprietary LLM (replace open-source Gemma)"),
        ("Ingestion", "Run full 9-stage pipeline on VZW requirements"),
        ("Tuning", "Stage-by-stage verification with corrections as feedback:\n"
         "    \u2022 Profile corrections (heading levels, req ID patterns)\n"
         "    \u2022 Taxonomy corrections (feature naming, missing features)\n"
         "    \u2022 Evaluation Q&A from domain experts (Excel)"),
        ("Validation", "Prove: use-case realization, accuracy targets (>80% completeness,\n"
         "    >90% accuracy, 100% citation quality), acceptable performance"),
    ]),
    ("Phase 2: Scale Out", "After PoC validation", ACCENT_BLUE, [
        ("Multi-MNO", "Add ATT, TMO document profiles and requirements"),
        ("Multi-Release", "Ingest multiple quarterly releases per MNO"),
        ("Cross-MNO queries", "Validate comparison and diff queries across carriers"),
        ("Test cases", "Implement test case parser (PoC Step 4), traceability queries"),
        ("Production LLM tuning", "Optimize prompts and context assembly for proprietary LLM"),
    ]),
    ("Phase 3: Production", "After scale-out validation", MID_BLUE, [
        ("Infrastructure", "Neo4j graph DB, production vector store (Milvus/Weaviate)"),
        ("Compliance Agent", "Automated compliance checking, delta sheet generation"),
        ("Integration", "REST API, CI/CD pipeline for automated ingestion on new releases"),
        ("Scale", "All MNOs \u00d7 all releases \u00d7 hundreds of documents"),
        ("Multimodal", "Diagram/call-flow understanding via multimodal LLM"),
    ]),
]

for i, (title, timeline, color, items) in enumerate(phases):
    x = Inches(0.4) + Inches(i * 4.25)
    y = Inches(1.5)

    header = add_shape(slide, x, y, Inches(4.0), Inches(0.8), color)
    add_text_box(slide, x + Inches(0.15), y + Inches(0.05), Inches(3.7), Inches(0.35),
                 title, font_size=16, color=WHITE, bold=True)
    add_text_box(slide, x + Inches(0.15), y + Inches(0.42), Inches(3.7), Inches(0.3),
                 timeline, font_size=11, color=RGBColor(0xDD, 0xDD, 0xDD))

    body_box = add_rounded_shape(slide, x, y + Inches(0.8), Inches(4.0), Inches(4.75),
                                 WHITE, color)

    txBox = slide.shapes.add_textbox(x + Inches(0.12), y + Inches(0.95),
                                     Inches(3.75), Inches(4.5))
    tf = txBox.text_frame
    tf.word_wrap = True

    for j, (label, desc) in enumerate(items):
        if j == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_after = Pt(6)

        r1 = p.add_run()
        r1.text = label + ": "
        r1.font.size = Pt(12)
        r1.font.color.rgb = color
        r1.font.bold = True
        r1.font.name = "Calibri"

        r2 = p.add_run()
        r2.text = desc
        r2.font.size = Pt(11)
        r2.font.color.rgb = DARK_GRAY
        r2.font.name = "Calibri"

    if i < len(phases) - 1:
        arrow = slide.shapes.add_shape(
            MSO_SHAPE.RIGHT_ARROW,
            x + Inches(4.02), Inches(3.5),
            Inches(0.22), Inches(0.4)
        )
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = MED_GRAY
        arrow.line.fill.background()

add_slide_number(slide, 13)


# ============================================================================
# SLIDE 14: Thank You / Q&A
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BLUE)

add_shape(slide, Inches(0), Inches(2.3), SLIDE_W, Inches(3.0), MID_BLUE)

add_text_box(slide, Inches(1), Inches(1.0), Inches(11), Inches(0.6),
             "NORA", font_size=48, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1), Inches(2.6), Inches(11), Inches(0.5),
             "Questions & Discussion", font_size=28,
             color=WHITE, alignment=PP_ALIGN.CENTER)

add_shape(slide, Inches(5.5), Inches(3.3), Inches(2.3), Pt(2), ACCENT_BLUE)

summary_items = [
    "Knowledge Graph + RAG hybrid solves cross-document, cross-MNO challenges that vanilla RAG cannot",
    "9-stage modular pipeline: each stage independently testable and tunable",
    "LLM-agnostic design: ready for proprietary LLM integration (Phase 1 next step)",
    "426 tests, working Web UI, metrics infrastructure \u2014 solid foundation for scale-out",
]

txBox = slide.shapes.add_textbox(Inches(2), Inches(3.7), Inches(9.3), Inches(2.5))
tf = txBox.text_frame
tf.word_wrap = True
for i, item in enumerate(summary_items):
    if i == 0:
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.space_after = Pt(8)
    r = p.add_run()
    r.text = "\u2022  " + item
    r.font.size = Pt(14)
    r.font.color.rgb = LIGHT_BLUE
    r.font.name = "Calibri"

add_text_box(slide, Inches(1), Inches(6.5), Inches(11), Inches(0.4),
             "github.com/kurnoolion/nora", font_size=13,
             color=RGBColor(0x88, 0xAA, 0xCC), alignment=PP_ALIGN.CENTER)


# ============================================================================
# Save
# ============================================================================
output_path = "NORA_Leadership_Presentation.pptx"
prs.save(output_path)
print(f"Presentation saved to {output_path}")
print(f"Total slides: {len(prs.slides)}")
